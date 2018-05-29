import argparse
import html
import json
import logging
import os
import re
from multiprocessing import Pool
from multiprocessing.util import Finalize

import pymysql
from pptx import Presentation
from tqdm import tqdm

from tqa import DEFAULTS, DATA_DIR
from tqa.retriever import utils
from tqa.retriever.tokenizer import CoreNlpTokenizer

log_format = logging.Formatter('%(asctime)s: [ %(message)s ]', '%Y/%m/%d %H:%M:%S')
console = logging.StreamHandler()
console.setFormatter(log_format)
logger = logging.getLogger()
logger.setLevel(logging.INFO)
logger.addHandler(console)

VALID_FORMAT = (".pptx", ".json", ".txt", ".srt")

TOKENIZER = None


def pool_init(tokenizer_heap):
    global TOKENIZER
    TOKENIZER = CoreNlpTokenizer(**{'heap': tokenizer_heap})
    Finalize(TOKENIZER, TOKENIZER.close, exitpriority=100)


def build(documents_dir, tokenizer_heap, num_workers=None):
    """ 预处理语料documents并存入sqlite
    :param: documents_dir: documents文件夹路径
    :param: num_workers: 读doc的并行数
    :return:
    """
    db_table = os.path.basename(documents_dir[:-1])

    conn = pymysql.connect(
        DEFAULTS['db_host'], DEFAULTS['db_user'], DEFAULTS['db_password'], DEFAULTS['db_database'],
        use_unicode=True, charset="utf8mb4"
    )
    cursor = conn.cursor()
    cursor.execute("SELECT COUNT(*) FROM information_schema.tables WHERE table_name = '%s'" % db_table)
    if cursor.fetchone()[0] == 0:
        logger.info('Creating database...')
        cursor.execute("""
          CREATE TABLE {table_name} (
            document_id VARCHAR(100) NOT NULL, 
            document_text MEDIUMTEXT,
            document_words MEDIUMTEXT,
            document_words_ws MEDIUMTEXT,
            document_span MEDIUMTEXT,
            document_pos MEDIUMTEXT,
            document_lemma MEDIUMTEXT,
            document_ner MEDIUMTEXT,
            PRIMARY KEY ( document_id ),
            INDEX [id_index] (document_id(100))
          ) DEFAULT CHARSET=utf8mb4
        """.format(table_name=db_table))
    else:
        logger.info('Database exists! Appending...')

    # 如果num_workers==None，则使用cpu_count()
    pool = Pool(num_workers, initializer=pool_init, initargs=(tokenizer_heap,))
    id_text_pairs = [id_text_pair for id_text_pair in documents_iterate(documents_dir, db_table, cursor)]

    count = 0
    # tqdm是一个快速，可扩展的进度条，可以在长循环中添加一个进度提示信息，只需要封装任意的迭代器tqdm(iterator)
    with tqdm(total=len(id_text_pairs)) as pbar:
        # imap(function, iterable)与map类似，对iterable中的每一项调用function
        for tuple in tqdm(pool.imap_unordered(tokenize, id_text_pairs, chunksize=1)):
            if tuple and tuple[2]:
                id = utils.normalize(tuple[0])
                text = utils.normalize(tuple[1])
                # logger.info((text[0:10] if len(text) > 10 else text) + "...")
                words = utils.normalize("\t".join(tuple[2].words()))
                words_ws = utils.normalize("\t".join(tuple[2].words_ws()))
                span = utils.normalize("\t".join([str(s[0]) + "," + str(s[1]) for s in tuple[2].span()]))
                pos = utils.normalize("\t".join(tuple[2].pos()))
                lemma = utils.normalize("\t".join(tuple[2].lemma()))
                ner = utils.normalize("\t".join(tuple[2].ner()))

                try:
                    cursor.execute("INSERT IGNORE INTO " + db_table + " VALUES (%s,%s,%s,%s,%s,%s,%s,%s)",
                                   (id, text, words, words_ws, span, pos, lemma, ner))
                    # conn.commit()
                    count += 1
                except pymysql.err.ProgrammingError as e:
                    conn.rollback()
                    print(e)
            pbar.update()
    conn.commit()
    logger.info('Committed total %d documents.' % count)
    conn.close()


def documents_iterate(documents_dir, db_table, cursor):
    """ 读取path下所有raw documents """
    if os.path.isdir(documents_dir):
        for dirpath, _, filenames in os.walk(documents_dir):
            for filename in filenames:
                if filename.endswith(VALID_FORMAT):
                    if filename.endswith(".pptx"):
                        document_path = os.path.join(dirpath, filename)
                        # ID为相对$RAW_DIR的文档路径（带后缀）
                        id = utils.normalize(document_path.replace(documents_dir, ''))
                        cursor.execute(
                            "SELECT COUNT(*) FROM %s WHERE document_id = '%s'" % (db_table, utils.normalize(id)))
                        if cursor.fetchone()[0] == 0:
                            presentation = Presentation(document_path)
                            text = handle_pptx(presentation)
                            yield (id, text)
                        else:
                            logger.info("IGNORE EXIST: " + id)
                    # PLAIN TEXT
                    elif filename.endswith(".txt"):
                        document_path = os.path.join(dirpath, filename)
                        id = utils.normalize(document_path.replace(documents_dir, ''))
                        cursor.execute(
                            "SELECT COUNT(*) FROM %s WHERE document_id = '%s'" % (db_table, utils.normalize(id)))
                        if cursor.fetchone()[0] == 0:
                            text = handle_txt(document_path)
                            yield (id, text)
                        else:
                            logger.info("IGNORE EXIST: " + id)
                    # SRT SUBTITLE
                    elif filename.endswith(".srt"):
                        document_path = os.path.join(dirpath, filename)
                        id = utils.normalize(document_path.replace(documents_dir, ''))
                        cursor.execute(
                            "SELECT COUNT(*) FROM %s WHERE document_id = '%s'" % (db_table, utils.normalize(id)))
                        if cursor.fetchone()[0] == 0:
                            text = handle_srt(document_path)
                            yield (id, text)
                        else:
                            logger.info("IGNORE EXIST: " + id)
                    # WIKI JSON
                    elif filename.endswith(".json"):
                        document_path = os.path.join(dirpath, filename)
                        with open(document_path) as file:
                            for line in file:
                                document = json.loads(line)
                                title = clean_title(document['title'])
                                id = utils.normalize(document['url'] + "@" + title)
                                text = document['text']
                                text = clean(text)
                                # cursor.execute("SELECT COUNT(*) FROM %s WHERE document_id = '%s'" % (db_table, id))
                                # if cursor.fetchone()[0] == 0:
                                yield (id, text)
                    else:
                        raise RuntimeError('Unsupported-format file: %s' % filename)
                else:
                    raise RuntimeError('Unsupported-format file: %s' % filename)
    else:
        raise RuntimeError('Path %s is not directory or invalid' % documents_dir)


def handle_srt(document_path):
    def clean_txt(t):
        t = re.sub(r"[\'\"\\]", " ", t)
        t = re.sub(u'\u3000', ' ', t)
        t = re.sub(r"\t+|\n+|\r+", "", t)  # 去除非空格的空白符
        t = re.sub(r"\s{2,}", " ", t)
        return t

    subtitles = []
    with open(document_path, encoding='utf-8') as file:
        all = file.readlines()
        for i in range(0, len(all)):
            if '-->' in all[i]:
                subtitles.append(all[i + 1])
                i += 1

        text = '\t'.join(subtitles)
        text = clean_txt(text)
        # 除去32bit的unicode，只保留16bit以下的unicode，否则在java分词的时候会出错
        text = "".join([char if ord(char) < 65535 else "?" for char in text])
        assert "  " not in text, "Contains double space!"
        return text


def handle_txt(document_path):
    def clean_txt(t):
        t = re.sub(r"[\'\"\\]", " ", t)
        t = re.sub(u'\u3000', ' ', t)
        t = re.sub(r"\t+|\n+|\r+", "", t)  # 去除非空格的空白符
        t = re.sub(r"\s{2,}", " ", t)
        return t

    with open(document_path, encoding='utf-8') as file:
        all = file.readlines()
        all = list(filter(lambda x: x and x.strip() != '', all))
        text = ' '.join(all)
        text = clean_txt(text)
        # 除去32bit的unicode，只保留16bit以下的unicode，否则在java分词的时候会出错
        text = "".join([char if ord(char) < 65535 else "?" for char in text])
        assert "  " not in text, "Contains double space!"
        return text


def handle_pptx(presentation):
    def clean_pptx(t):
        t = re.sub(r"[\'\"\\]", " ", t)
        t = re.sub(r"\t+|\n+|\r+", "", t)  # 去除非空格的空白符
        t = re.sub(r"\s{2,}", " ", t)
        t = re.sub(r"[.。!！?？]", "；", t)
        return t

    slides_text = []
    first = True
    for slide in presentation.slides:
        if first:
            first = False
            continue
        title = ''
        if hasattr(slide.shapes.title, "text"):
            title = slide.shapes.title.text.strip()
            if title and title.strip() != '':
                title = re.sub(r"\n+|\r+", " ", title)
                title = clean_pptx(title)
                title = title + "："
        slide_text = []
        start = 0 if title == '' else 1
        end = len(slide.shapes)
        if start < end:
            for i in range(start, end):
                try:
                    shape = slide.shapes[i]
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text and text.strip() != '':
                            text = re.sub(r"\n+|\r+", "；", text)
                            text = clean_pptx(text)
                            slide_text.append(text)
                    if hasattr(shape, "image"):
                        slide_text.append("（图见PPT）")
                except KeyError as e:
                    pass
        content = title + "；".join(slide_text)
        if content and content.strip() != '':
            slides_text.append(content)
    return '。'.join(slides_text)


def tokenize(id_text_pairs):
    """ 转换raw文件为text并存入db """
    global TOKENIZER
    id = id_text_pairs[0]
    text = id_text_pairs[1]
    tokens = TOKENIZER.tokenize(text) if text else None
    return utils.normalize(id), utils.normalize(text), tokens


def clean_title(title):
    title = title.strip()
    title = re.sub(r"[\'\"\\]", " ", title)
    title = re.sub(r"\t+|\n+|\r+", "", title)  # 去除非空格的空白符
    title = re.sub(r"\s{2,}", "", title)  # 去除2至多个空格符
    # 去除'-{zh-hans: 高管; zh-hant: 高阶主管}-', '-{zh-hant:战网;zh-hans:Battle.net;}-'
    title = re.sub(r'-{.*?zh-hans:\s?(.*?)\s?;.*?}-', r'\1', title)
    # 去除'-{zh-hant:痴呆症;zh-hans:失智症}-'
    title = re.sub(r'-{.*?zh-hans:\s?(.*?)\s?}-', r'\1', title)
    # 去除'-{zh-cn:特奥伊奴·特奥;zh-hk:迪安尼·泰希奥}-'
    title = re.sub(r'-{.*?zh-cn:\s?(.*?)\s?;.*?}-', r'\1', title)
    # 去除'-{zh-tw:线;zh-hk:线;zh-cn:线}-'
    title = re.sub(r'-{.*?zh-cn:\s?(.*?)\s?}-', r'\1', title)
    # 去除'-{里}-', '-{}-'
    title = re.sub(r'-{(.*?)}-', r'\1', title)
    # 去除'（）'
    title = re.sub(r'（\s*?）', "", title)
    title = re.sub(r'\(\s*?\)', "", title)
    title = re.sub(r" ", "_", title)
    title = re.sub(r"@", "_", title)
    return title


def clean(text):
    # 格式化text
    text = html.unescape(text)  # 去除html标签
    text = re.sub(r"[\'\"\\]", "", text)
    text = re.sub(r'^.*?\s+', '', text)  # 去除标题
    text = re.sub(r"\t+|\n+|\r+", "", text)  # 去除非空格的空白符
    text = re.sub(r"\s{2,}", "", text)  # 去除2至多个空格符
    text = re.sub(r'（.*?），?是?[^。]', '是', text, 1)  # 去除第一个（）
    text = re.sub(r'，是', '是', text)  # 去除，
    # 去除'-{zh-hans: 高管; zh-hant: 高阶主管}-', '-{zh-hant:战网;zh-hans:Battle.net;}-'
    text = re.sub(r'-{.*?zh-hans:\s?(.*?)\s?;.*?}-', r'\1', text)
    # 去除'-{zh-hant:痴呆症;zh-hans:失智症}-'
    text = re.sub(r'-{.*?zh-hans:\s?(.*?)\s?}-', r'\1', text)
    # 去除'-{zh-cn:特奥伊奴·特奥;zh-hk:迪安尼·泰希奥}-'
    text = re.sub(r'-{.*?zh-cn:\s?(.*?)\s?;.*?}-', r'\1', text)
    # 去除'-{zh-tw:线;zh-hk:线;zh-cn:线}-'
    text = re.sub(r'-{.*?zh-cn:\s?(.*?)\s?}-', r'\1', text)
    # 去除'-{里}-', '-{}-'
    text = re.sub(r'-{(.*?)}-', r'\1', text)
    # 去除'（）'
    text = re.sub(r'（\s*?）', "", text)
    text = re.sub(r'\(\s*?\)', "", text)

    text = re.sub(r"\s{2,}", "", text)  # 去除2至多个空格符

    # 除去32bit的unicode，只保留16bit以下的unicode，否则在java分词的时候会出错
    text = "".join([char if ord(char) < 65535 else "?" for char in text])
    assert "  " not in text, "Contains double space!"
    return utils.normalize(text)


if __name__ == '__main__':
    parser = argparse.ArgumentParser('db_builder.py', formatter_class=argparse.ArgumentDefaultsHelpFormatter)
    parser.add_argument('--documents-dir', type=str, default=None)
    parser.add_argument('--tokenizer-heap', type=str, default=DEFAULTS['tokenizer_heap'])
    parser.add_argument('--num-workers', type=int, default=DEFAULTS['num_workers'])
    args = parser.parse_args()

    args.documents_dir = os.path.join(DATA_DIR, args.documents_dir)
    build(
        documents_dir=args.documents_dir,
        tokenizer_heap=args.tokenizer_heap,
        num_workers=args.num_workers
    )
